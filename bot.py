import logging
import os
import json
import base64
import re
import sys
import threading
import time
from datetime import datetime
from collections import defaultdict

import psutil
import humanize
import requests
from flask import Flask
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    Application, CommandHandler, MessageHandler,
    CallbackQueryHandler, filters, ContextTypes
)

# ========== Keys ==========
BOT_TOKEN = os.environ.get('BOT_TOKEN')
OPENROUTER_KEY = os.environ.get('OPENROUTER_KEY')

if not BOT_TOKEN:
    raise ValueError("❌ BOT_TOKEN غير موجود في Environment Variables!")

# ========== Flask Web Server (single instance) ==========
flask_app = Flask(__name__)

@flask_app.route('/')
def home():
    return "✅ البوت يعمل بنجاح!"

@flask_app.route('/health')
def health():
    return "OK", 200

_webserver_started = False

def run_webserver():
    global _webserver_started
    if _webserver_started:
        return
    _webserver_started = True
    port = int(os.environ.get('PORT', 5000))
    try:
        flask_app.run(host='0.0.0.0', port=port, debug=False, use_reloader=False)
    except Exception as e:
        print(f"⚠️ خطأ في خادم الويب: {e}")

threading.Thread(target=run_webserver, daemon=True).start()
time.sleep(1)
print(f"🌐 خادم الويب نشط على المنفذ {os.environ.get('PORT', 5000)}")

# ========== Logging ==========
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)

# ========== IDs ==========
# تُقرأ من الأسرار أولاً — وإلا تستخدم القيمة الافتراضية
OWNER_ID        = int(os.environ.get('OWNER_ID',        383022213))
ABDULKHALIQ_ID  = int(os.environ.get('ABDULKHALIQ_ID',  6818088581))
FATIMA_ID       = int(os.environ.get('FATIMA_ID',       1295746334))

# مجموعات الصلاحيات
FULL_OWNERS    = {OWNER_ID, ABDULKHALIQ_ID}          # صلاحيات كاملة كالمالك تماماً
EXAM_PARTNERS  = {OWNER_ID, ABDULKHALIQ_ID, FATIMA_ID}  # أوامر الاختبارات فقط

# ========== Files & Dirs ==========
HISTORY_FILE = "conversations.json"
USERS_COUNT_FILE = "users_count.json"
QUIZZES_DIR = "quizzes"
os.makedirs(QUIZZES_DIR, exist_ok=True)

# ========== User Counter ==========
def load_user_count():
    try:
        if os.path.exists(USERS_COUNT_FILE):
            with open(USERS_COUNT_FILE, 'r') as f:
                data = json.load(f)
                return data.get('count', 567), set(data.get('seen_ids', []))
    except Exception:
        pass
    return 567, set()

def save_user_count(count, seen_ids):
    try:
        with open(USERS_COUNT_FILE, 'w') as f:
            json.dump({'count': count, 'seen_ids': list(seen_ids)}, f)
    except Exception:
        pass

user_count, seen_user_ids = load_user_count()

needs_description_update = False

def track_user(user_id):
    global user_count, seen_user_ids, needs_description_update
    if user_id not in seen_user_ids and user_id not in FULL_OWNERS:
        seen_user_ids.add(user_id)
        user_count += 1
        save_user_count(user_count, seen_user_ids)
        needs_description_update = True
        return True
    return False

async def update_bot_description(context: ContextTypes.DEFAULT_TYPE):
    global needs_description_update
    try:
        description = (
            f"🎓 بوت اختبارات الجامعة بالذكاء الاصطناعي\n"
            f"━━━━━━━━━━━━━━━━━━━━\n"
            f"👥 المستخدمون: {user_count:,} طالب\n"
            f"━━━━━━━━━━━━━━━━━━━━\n"
            f"📤 أرسل ملف PDF أو PPTX لتوليد اختبار\n"
            f"📝 /quiz لحل الاختبارات المتاحة\n"
            f"🤖 دردشة ذكية باللغة العربية"
        )
        await context.bot.set_my_description(description)
        needs_description_update = False
    except Exception as e:
        logging.warning(f"تحديث الوصف: {e}")

# ========== In-Memory Data ==========
users_db = {}
conversation_history = defaultdict(list)
performance_stats = {
    'total_messages_processed': 0,
    'total_api_calls': 0,
    'total_tokens_estimated': 0,
    'start_time': datetime.now(),
    'peak_memory': 0
}

# quiz_sessions: {user_id: {quiz_id, current_q, score, answers}}
quiz_sessions = {}

# pending_quiz: owner uploaded file, waiting for instructions
pending_quiz = {}

# ========== Load Conversations ==========
try:
    if os.path.exists(HISTORY_FILE):
        file_size = os.path.getsize(HISTORY_FILE)
        with open(HISTORY_FILE, 'r', encoding='utf-8') as f:
            loaded_history = json.load(f)
            for key, value in loaded_history.items():
                conversation_history[int(key)] = value
        logging.info(f"✅ تم تحميل المحادثات (الحجم: {humanize.naturalsize(file_size)})")
except Exception as e:
    logging.error(f"❌ فشل تحميل المحادثات: {e}")

# ========== Memory & Stats Helpers ==========
def get_memory_usage():
    process = psutil.Process(os.getpid())
    memory_info = process.memory_info()
    rss = memory_info.rss
    vms = memory_info.vms
    global performance_stats
    if rss > performance_stats['peak_memory']:
        performance_stats['peak_memory'] = rss
    return {
        'rss': rss, 'vms': vms,
        'rss_human': humanize.naturalsize(rss),
        'vms_human': humanize.naturalsize(vms),
        'percent': process.memory_percent(),
        'cpu_percent': process.cpu_percent()
    }

def estimate_conversation_size():
    total_chars = sum(len(msg.get('content', '')) for h in conversation_history.values() for msg in h)
    total_messages = sum(len(h) for h in conversation_history.values())
    estimated_bytes = total_chars * 2
    return {
        'total_messages': total_messages,
        'total_chars': total_chars,
        'estimated_bytes': estimated_bytes,
        'estimated_human': humanize.naturalsize(estimated_bytes),
        'users_count': len(conversation_history)
    }

def save_conversations():
    try:
        to_save = {str(k): v for k, v in conversation_history.items()}
        with open(HISTORY_FILE, 'w', encoding='utf-8') as f:
            json.dump(to_save, f, ensure_ascii=False, indent=2)
        logging.info(f"✅ تم حفظ المحادثات")
        return True
    except Exception as e:
        logging.error(f"❌ فشل حفظ المحادثات: {e}")
        return False

# ========== AI API Helpers ==========
def call_ai_g4f(messages):
    """استدعاء الذكاء الاصطناعي عبر g4f مجاناً"""
    try:
        from g4f.client import Client
        from g4f.Provider import PollinationsAI, Copilot
        client = Client(provider=PollinationsAI)
        response = client.chat.completions.create(
            model="openai-fast",
            messages=messages
        )
        return response.choices[0].message.content
    except Exception as e:
        logging.warning(f"g4f PollinationsAI فشل: {e} — جاري المحاولة مع Copilot")
        try:
            from g4f.client import Client
            from g4f.Provider import Copilot
            client = Client(provider=Copilot)
            response = client.chat.completions.create(
                model="gpt-4",
                messages=messages
            )
            return response.choices[0].message.content
        except Exception as e2:
            logging.warning(f"g4f Copilot فشل: {e2}")
            raise Exception(f"g4f فشل: {e2}")

def call_ai_openrouter(messages, model="meta-llama/llama-3-8b-instruct"):
    """استدعاء OpenRouter كبديل احتياطي"""
    if not OPENROUTER_KEY:
        raise Exception("OPENROUTER_KEY غير متوفر، g4f فشل ولا يوجد بديل")
    response = requests.post(
        url="https://openrouter.ai/api/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {OPENROUTER_KEY}",
            "Content-Type": "application/json",
        },
        json={"model": model, "messages": messages, "temperature": 0.3},
        timeout=120
    )
    data = response.json()
    if response.status_code == 200:
        return data['choices'][0]['message']['content']
    raise Exception(data.get('error', {}).get('message', 'خطأ في OpenRouter'))

def call_ai(messages, model="meta-llama/llama-3-8b-instruct"):
    """استدعاء الذكاء الاصطناعي — يجرب g4f أولاً ثم OpenRouter كبديل"""
    performance_stats['total_api_calls'] += 1
    try:
        result = call_ai_g4f(messages)
        logging.info("✅ g4f نجح")
        return result
    except Exception as e:
        logging.warning(f"g4f فشل، يتم التحويل لـ OpenRouter: {e}")
        return call_ai_openrouter(messages, model)

def call_vision_ai(image_base64, prompt, mime_type="image/jpeg"):
    messages = [{
        "role": "user",
        "content": [
            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_base64}"}},
            {"type": "text", "text": prompt}
        ]
    }]
    response = requests.post(
        url="https://openrouter.ai/api/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {OPENROUTER_KEY}",
            "Content-Type": "application/json",
        },
        json={"model": "google/gemini-flash-1.5", "messages": messages, "temperature": 0.3},
        timeout=120
    )
    performance_stats['total_api_calls'] += 1
    data = response.json()
    if response.status_code == 200:
        return data['choices'][0]['message']['content']
    raise Exception(data.get('error', {}).get('message', 'خطأ في Vision API'))

# ========== Google Drive Integration ==========
GOOGLE_API_KEY = os.environ.get('GOOGLE_API_KEY')
DRIVE_ROOT_FOLDER_ID = "1NbqIyWzQTLTRXj1vYTADSvcSSgUo-0NO"

drive_sessions = {}

def drive_list(folder_id):
    """قائمة محتويات مجلد Google Drive"""
    try:
        url = "https://www.googleapis.com/drive/v3/files"
        params = {
            "q": f"'{folder_id}' in parents and trashed=false",
            "key": GOOGLE_API_KEY,
            "fields": "files(id,name,mimeType)",
            "pageSize": 50,
            "orderBy": "name"
        }
        r = requests.get(url, params=params, timeout=15)
        if r.status_code == 200:
            return r.json().get('files', [])
        return []
    except Exception as e:
        logging.error(f"Drive list error: {e}")
        return []

def drive_download(file_id, dest_path):
    """تحميل ملف من Google Drive"""
    try:
        url = f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media&key={GOOGLE_API_KEY}"
        r = requests.get(url, timeout=120, stream=True)
        if r.status_code == 200:
            with open(dest_path, 'wb') as f:
                for chunk in r.iter_content(chunk_size=8192):
                    f.write(chunk)
            return True
        return False
    except Exception as e:
        logging.error(f"Drive download error: {e}")
        return False

async def drive_book_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """أمر /كتاب — تصفح كتب Google Drive"""
    user_id = update.effective_user.id
    if user_id not in EXAM_PARTNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط.")
        return
    folders = drive_list(DRIVE_ROOT_FOLDER_ID)
    if not folders:
        await update.message.reply_text("❌ تعذر الوصول إلى مجلد الكتب.")
        return
    keyboard = []
    for f in folders:
        if f['mimeType'] == 'application/vnd.google-apps.folder':
            keyboard.append([InlineKeyboardButton(
                f"📚 {f['name']}", callback_data=f"drv_sub_{f['id']}_{f['name']}"
            )])
    await update.message.reply_text(
        "📖 *اختر المادة:*",
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode='Markdown'
    )

async def handle_drive_callback(query, user_id, data):
    """معالجة callbacks الخاصة بـ Google Drive"""
    parts = data.split("_", 3)

    if data.startswith("drv_sub_"):
        folder_id = parts[2]
        folder_name = parts[3] if len(parts) > 3 else "المادة"
        items = drive_list(folder_id)
        keyboard = []
        for item in items:
            if item['mimeType'] == 'application/vnd.google-apps.folder':
                keyboard.append([InlineKeyboardButton(
                    f"📁 {item['name']}", callback_data=f"drv_grade_{item['id']}_{item['name']}"
                )])
            elif 'pdf' in item['mimeType'].lower() or item['name'].endswith('.pdf'):
                keyboard.append([InlineKeyboardButton(
                    f"📄 {item['name']}", callback_data=f"drv_file_{item['id']}_{item['name'][:40]}"
                )])
        keyboard.append([InlineKeyboardButton("🔙 رجوع", callback_data="drv_back_root")])
        await query.message.reply_text(
            f"📚 *{folder_name}* — اختر الصف أو الملف:",
            reply_markup=InlineKeyboardMarkup(keyboard),
            parse_mode='Markdown'
        )

    elif data.startswith("drv_grade_"):
        folder_id = parts[2]
        folder_name = parts[3] if len(parts) > 3 else "الصف"
        items = drive_list(folder_id)
        keyboard = []
        for item in items:
            if item['mimeType'] == 'application/vnd.google-apps.folder':
                keyboard.append([InlineKeyboardButton(
                    f"📁 {item['name']}", callback_data=f"drv_grade_{item['id']}_{item['name']}"
                )])
            elif 'pdf' in item['mimeType'].lower() or item['name'].endswith('.pdf'):
                keyboard.append([InlineKeyboardButton(
                    f"📄 {item['name']}", callback_data=f"drv_file_{item['id']}_{item['name'][:40]}"
                )])
        keyboard.append([InlineKeyboardButton("🔙 رجوع", callback_data="drv_back_root")])
        await query.message.reply_text(
            f"📁 *{folder_name}* — اختر الملف:",
            reply_markup=InlineKeyboardMarkup(keyboard),
            parse_mode='Markdown'
        )

    elif data.startswith("drv_file_"):
        file_id = parts[2]
        file_name = parts[3] if len(parts) > 3 else "الكتاب"
        drive_sessions[user_id] = {'file_id': file_id, 'file_name': file_name, 'step': 'ask_count'}
        await query.message.reply_text(
            f"📄 *{file_name}*\n\n"
            "كم عدد الأسئلة التي تريد توليدها؟\n"
            "أرسل رقماً من 5 إلى 200",
            parse_mode='Markdown'
        )

    elif data == "drv_back_root":
        folders = drive_list(DRIVE_ROOT_FOLDER_ID)
        keyboard = []
        for f in folders:
            if f['mimeType'] == 'application/vnd.google-apps.folder':
                keyboard.append([InlineKeyboardButton(
                    f"📚 {f['name']}", callback_data=f"drv_sub_{f['id']}_{f['name']}"
                )])
        await query.message.reply_text(
            "📖 *اختر المادة:*",
            reply_markup=InlineKeyboardMarkup(keyboard),
            parse_mode='Markdown'
        )

async def process_drive_file(update, context, user_id, file_id, file_name, question_count):
    """تحميل ملف من Drive وتوليد الأسئلة منه"""
    await update.message.reply_text(
        f"⬇️ جاري تحميل *{file_name}* من Drive...",
        parse_mode='Markdown'
    )
    tmp_path = f"/tmp/drive_{user_id}.pdf"
    success = drive_download(file_id, tmp_path)
    if not success:
        await update.message.reply_text("❌ فشل تحميل الملف. تأكد من صلاحيات المشاركة.")
        return
    await update.message.reply_text("📖 جاري قراءة الكتاب...")
    try:
        import fitz
        doc = fitz.open(tmp_path)
        text = "".join(page.get_text() for page in doc)
        doc.close()
        os.remove(tmp_path)
    except Exception as e:
        await update.message.reply_text(f"❌ خطأ في قراءة الملف: {e}")
        return
    if not text or len(text.strip()) < 100:
        await update.message.reply_text("❌ لم يتم استخراج نص كافٍ من الملف.")
        return
    instructions = f"{question_count} سؤال من كتاب {file_name}"
    await create_quiz_from_content(update, context, text, None, instructions)

# ========== File Extraction ==========
def extract_pdf_text(file_path):
    try:
        import fitz
        doc = fitz.open(file_path)
        text = "".join(page.get_text() for page in doc)
        doc.close()
        return text
    except Exception as e:
        logging.error(f"PDF error: {e}")
        return None

def extract_pptx_text(file_path):
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logging.error(f"PPTX error: {e}")
        return None

# ========== Quiz Storage ==========
def get_quiz_list():
    quizzes = []
    for f in sorted(os.listdir(QUIZZES_DIR)):
        if f.endswith('.json'):
            try:
                with open(os.path.join(QUIZZES_DIR, f), 'r', encoding='utf-8') as fp:
                    data = json.load(fp)
                    quizzes.append({
                        'id': f[:-5],
                        'title': data.get('title', f[:-5]),
                        'questions_count': len(data.get('questions', []))
                    })
            except Exception:
                pass
    return quizzes

def load_quiz(quiz_id):
    path = os.path.join(QUIZZES_DIR, f"{quiz_id}.json")
    if os.path.exists(path):
        with open(path, 'r', encoding='utf-8') as f:
            return json.load(f)
    return None

def save_quiz(quiz_id, quiz_data):
    path = os.path.join(QUIZZES_DIR, f"{quiz_id}.json")
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(quiz_data, f, ensure_ascii=False, indent=2)

# ========== Quiz Generation ==========
BATCH_SIZE = 20

def parse_question_count(instructions):
    """استخراج عدد الأسئلة المطلوبة من التعليمات"""
    match = re.search(r'(\d+)\s*سؤال', instructions)
    if match:
        return int(match.group(1))
    match = re.search(r'(\d+)', instructions)
    if match:
        n = int(match.group(1))
        if 1 <= n <= 500:
            return n
    return 10

def parse_quiz_title(instructions, text):
    """استخراج عنوان مقترح للاختبار من أول سطور النص"""
    first_line = text.strip().split('\n')[0][:60] if text else "اختبار"
    return first_line if first_line else "اختبار جديد"

CHUNK_SIZE = 40000

def get_text_chunk(text, batch_num, total_batches):
    """تقسيم النص على الدفعات حتى يغطي كل دفعة جزءاً مختلفاً من الكتاب"""
    total_len = len(text)
    chunk_len = min(CHUNK_SIZE, total_len)
    if total_batches <= 1:
        return text[:chunk_len]
    step = max(1, (total_len - chunk_len) // (total_batches - 1))
    start = step * (batch_num - 1)
    return text[start:start + chunk_len]

def generate_batch(text, instructions, batch_num, total_batches, batch_size, already_generated):
    """توليد دفعة واحدة من الأسئلة"""
    already_note = ""
    if already_generated > 0:
        already_note = f"لقد أنشأت {already_generated} سؤالاً بالفعل. "

    text_chunk = get_text_chunk(text, batch_num, total_batches)

    prompt = f"""أنت خبير متخصص في إنشاء الاختبارات الأكاديمية.
{already_note}أنشئ {batch_size} سؤالاً جديداً (الدفعة {batch_num} من {total_batches}) بناءً على النص أدناه.
التعليمات: {instructions}

النص:
{text_chunk}

مهم جداً:
- لا تكرر الأسئلة السابقة
- غطّ أجزاء مختلفة من النص
- أعطني النتيجة بصيغة JSON فقط بدون أي نص إضافي أو markdown

الصيغة المطلوبة (JSON فقط، مصفوفة questions فقط):
[
  {{
    "question": "نص السؤال",
    "type": "multiple_choice",
    "options": ["أ. الخيار الأول", "ب. الخيار الثاني", "ج. الخيار الثالث", "د. الخيار الرابع"],
    "correct": "أ",
    "explanation": "شرح الإجابة الصحيحة"
  }},
  {{
    "question": "نص السؤال الصح/خطأ",
    "type": "true_false",
    "options": ["✅ صح", "❌ خطأ"],
    "correct": "✅ صح",
    "explanation": "شرح الإجابة"
  }}
]"""

    result = call_ai(
        [{"role": "user", "content": prompt}],
        model="anthropic/claude-3-haiku"
    )
    result = result.strip()
    result = re.sub(r'```json\s*', '', result)
    result = re.sub(r'```\s*', '', result)
    result = result.strip()
    parsed = json.loads(result)
    if isinstance(parsed, list):
        return parsed
    if isinstance(parsed, dict) and 'questions' in parsed:
        return parsed['questions']
    raise ValueError("تنسيق غير متوقع من AI")

# ========== Notification Helpers ==========
async def send_to_owner(context, text):
    try:
        await context.bot.send_message(chat_id=OWNER_ID, text=text, parse_mode='Markdown')
    except Exception as e:
        logging.error(f"فشل إرسال للمالك: {e}")

async def send_to_user(context, user_id, text):
    try:
        await context.bot.send_message(chat_id=user_id, text=text, parse_mode='Markdown')
        return True
    except Exception as e:
        logging.error(f"فشل إرسال للمستخدم {user_id}: {e}")
        return False

# ========== Quiz UI Helpers ==========
async def show_quiz_list(update, context):
    quizzes = get_quiz_list()
    if not quizzes:
        await update.message.reply_text(
            "📚 لا توجد اختبارات متاحة حالياً.\n\n"
            "سيتم إضافة اختبارات قريباً! 🔜"
        )
        return
    keyboard = []
    for q in quizzes:
        keyboard.append([InlineKeyboardButton(
            f"📝 {q['title']} ({q['questions_count']} سؤال)",
            callback_data=f"start_quiz_{q['id']}"
        )])
    await update.message.reply_text(
        f"📚 *الاختبارات المتاحة*\n\n"
        f"👥 عدد المستخدمين: *{user_count}*\n\n"
        "اختر اختباراً للبدء:",
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode='Markdown'
    )

async def send_question(send_target, context, user_id, quiz_data, question_index):
    questions = quiz_data['questions']
    if question_index >= len(questions):
        await show_quiz_result(send_target, context, user_id, quiz_data)
        return
    q = questions[question_index]
    total = len(questions)
    keyboard = []
    for i, option in enumerate(q['options']):
        keyboard.append([InlineKeyboardButton(
            option, callback_data=f"answer_{question_index}_{i}"
        )])
    text = f"❓ *سؤال {question_index + 1} من {total}*\n\n{q['question']}"
    await send_target.reply_text(text, reply_markup=InlineKeyboardMarkup(keyboard), parse_mode='Markdown')

async def show_quiz_result(send_target, context, user_id, quiz_data):
    session = quiz_sessions.get(user_id, {})
    score = session.get('score', 0)
    total = len(quiz_data['questions'])
    percentage = (score / total * 100) if total > 0 else 0

    if percentage >= 90:
        rating = "🏆 ممتاز!"
    elif percentage >= 75:
        rating = "⭐ جيد جداً"
    elif percentage >= 60:
        rating = "👍 مقبول"
    else:
        rating = "💪 حاول مرة أخرى"

    quiz_id = session.get('quiz_id', '')
    keyboard = [[
        InlineKeyboardButton("🔄 إعادة الاختبار", callback_data=f"start_quiz_{quiz_id}"),
        InlineKeyboardButton("📚 اختبار آخر", callback_data="show_quizzes")
    ]]
    await send_target.reply_text(
        f"🎯 *انتهى الاختبار!*\n\n"
        f"📊 *نتيجتك:* {score} / {total}\n"
        f"📈 *النسبة:* {percentage:.0f}%\n"
        f"*التقييم:* {rating}\n\n"
        "شكراً على مشاركتك! 🙏",
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode='Markdown'
    )
    quiz_sessions.pop(user_id, None)

async def create_quiz_from_content(update, context, text_content, image_base64, instructions):
    try:
        if image_base64 and not text_content:
            await update.message.reply_text("🔍 جاري قراءة الصورة بالذكاء الاصطناعي...")
            extract_prompt = "استخرج جميع النصوص والمعلومات من هذه الصورة بالتفصيل الكامل."
            text_content = call_vision_ai(image_base64, extract_prompt)

        if not text_content or len(text_content.strip()) < 50:
            await update.message.reply_text("❌ المحتوى المستخرج قصير جداً. تأكد من الملف.")
            return

        total_questions = parse_question_count(instructions)
        total_questions = min(total_questions, 200)

        num_batches = (total_questions + BATCH_SIZE - 1) // BATCH_SIZE
        quiz_title = parse_quiz_title(instructions, text_content)

        if num_batches == 1:
            await update.message.reply_text(
                f"🤖 جاري توليد *{total_questions}* سؤالاً...",
                parse_mode='Markdown'
            )
        else:
            await update.message.reply_text(
                f"🤖 سيتم توليد *{total_questions}* سؤالاً على *{num_batches}* دفعات ({BATCH_SIZE} سؤال لكل دفعة)\n\n"
                f"⏳ جاري البدء...",
                parse_mode='Markdown'
            )

        all_questions = []
        quiz_id = f"quiz_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        progress_msg = None

        for batch_num in range(1, num_batches + 1):
            remaining = total_questions - len(all_questions)
            batch_size = min(BATCH_SIZE, remaining)
            if batch_size <= 0:
                break

            try:
                if num_batches > 1:
                    progress_text = (
                        f"⚙️ *جاري التوليد...*\n\n"
                        f"📦 الدفعة {batch_num} من {num_batches}\n"
                        f"✅ تم إنشاء: {len(all_questions)} سؤال\n"
                        f"🎯 المتبقي: {total_questions - len(all_questions)} سؤال\n\n"
                        + "▓" * batch_num + "░" * (num_batches - batch_num) + f" {int(batch_num/num_batches*100)}%"
                    )
                    if progress_msg is None:
                        progress_msg = await update.message.reply_text(progress_text, parse_mode='Markdown')
                    else:
                        try:
                            await progress_msg.edit_text(progress_text, parse_mode='Markdown')
                        except Exception:
                            pass

                batch_questions = generate_batch(
                    text_content, instructions,
                    batch_num, num_batches, batch_size,
                    len(all_questions)
                )
                all_questions.extend(batch_questions)
                logging.info(f"✅ دفعة {batch_num}: تم توليد {len(batch_questions)} سؤال")

                quiz_data = {
                    'title': quiz_title,
                    'questions': all_questions
                }
                save_quiz(quiz_id, quiz_data)

            except Exception as e:
                logging.error(f"خطأ في الدفعة {batch_num}: {e}")
                if len(all_questions) > 0:
                    await update.message.reply_text(
                        f"⚠️ خطأ في الدفعة {batch_num}: {str(e)[:100]}\n"
                        f"سيتم حفظ الأسئلة المولّدة حتى الآن ({len(all_questions)} سؤال)."
                    )
                    break
                else:
                    raise

        if not all_questions:
            await update.message.reply_text("❌ لم يتم توليد أي أسئلة. حاول مرة أخرى.")
            return

        final_quiz_data = {'title': quiz_title, 'questions': all_questions}
        save_quiz(quiz_id, final_quiz_data)

        if progress_msg:
            try:
                await progress_msg.edit_text(
                    f"✅ *اكتمل التوليد!*\n\n"
                    "▓" * num_batches + f" 100%",
                    parse_mode='Markdown'
                )
            except Exception:
                pass

        keyboard = [[InlineKeyboardButton("▶️ ابدأ الاختبار الآن", callback_data=f"start_quiz_{quiz_id}")]]
        await update.message.reply_text(
            f"✅ *تم إنشاء الاختبار بنجاح!*\n\n"
            f"📝 *العنوان:* {quiz_title}\n"
            f"❓ *عدد الأسئلة:* {len(all_questions)}\n"
            f"📦 *عدد الدفعات:* {num_batches}\n\n"
            f"الاختبار متاح الآن للمستخدمين عبر `/quiz`",
            reply_markup=InlineKeyboardMarkup(keyboard),
            parse_mode='Markdown'
        )

    except json.JSONDecodeError:
        await update.message.reply_text("❌ خطأ في تنسيق رد AI. حاول مرة أخرى مع تعليمات أوضح.")
    except Exception as e:
        await update.message.reply_text(f"❌ خطأ في توليد الاختبار: {str(e)[:200]}")
        logging.error(f"Quiz generation error: {e}")

# ========== Command Handlers ==========
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    user_name = update.message.from_user.first_name
    username = update.message.from_user.username
    is_new = track_user(user_id)

    users_db[user_id] = {
        'name': user_name,
        'username': username,
        'first_seen': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    }

    if user_id == OWNER_ID:
        memory = get_memory_usage()
        conv_stats = estimate_conversation_size()
        stats = get_platform_stats()
        welcome_text = (
            f"👑 *مرحباً أيها المالك!*\n\n"
            f"👥 المستخدمين: {user_count} | طلاب المنصة: {stats.get('students', 0)}\n"
            f"📚 أسئلة المنصة: {stats.get('questions', 0)} | 💾 {memory['rss_human']}\n\n"
            "🏫 *منصة رواد التحصيلي:*\n"
            "• `/platform` — قائمة المنصة الكاملة\n"
            "• `/daily physics|chemistry|biology|math` — تدريب يومي\n"
            "• `/assessment physics|...` — اختبار تقييم\n"
            "• `/ch1 physics|...` — فصل أول\n"
            "• `/ch2 physics|...` — فصل ثاني\n\n"
            "🔴 *اختبار مباشر (في المجموعة):*\n"
            "• `/live physics 10` — ابدأ اختباراً حياً\n"
            "• `/next` — السؤال التالي\n"
            "• `/leaderboard` — الترتيب الحالي\n"
            "• `/endlive` — إنهاء الاختبار\n\n"
            "📝 *أوامر إضافية:*\n"
            "• `/quiz` `/quizzes` `/users` `/stats` `/memory`\n"
            "• `/send` `/broadcast` `/clear` `/save`\n\n"
            "📎 أرسل PDF/TXT/PPTX/صورة لإنشاء اختبار"
        )
    elif user_id == FATIMA_ID:
        welcome_text = f"🌸 *أهلاً فاطمة المطيري!* 🌸\n\n👥 عدد المستخدمين: {user_count}"
        await send_to_owner(context, f"🌟 فاطمة المطيري دخلت البوت")
    elif user_id == ABDULKHALIQ_ID:
        welcome_text = f"👋 *مرحباً عبدالخالق!*\n\n👥 عدد المستخدمين: {user_count}"
        await send_to_owner(context, f"👤 عبدالخالق دخل البوت")
    else:
        welcome_text = (
            f"👋 *مرحباً {user_name}!*\n\n"
            f"أنا بوت الذكاء الاصطناعي للاختبارات الجامعية 🎓\n\n"
            f"👥 عدد المستخدمين: *{user_count}*\n\n"
            "✨ *يمكنني:*\n"
            "• الإجابة على أسئلتك بالعربية\n"
            "• `/quiz` - حل اختبارات تفاعلية\n\n"
            "أرسل لي أي سؤال وسأجيبك! 🤖"
        )
        if is_new:
            user_info = f"@{username}" if username else f"معرف {user_id}"
            await send_to_owner(
                context,
                f"🆕 *مستخدم جديد* (#{user_count})\n"
                f"👤 {user_name}\n🆔 `{user_id}`\n📱 {user_info}"
            )

    await update.message.reply_text(welcome_text, parse_mode='Markdown')

async def quiz_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    track_user(user_id)
    await show_quiz_list(update, context)

async def quizzes_manage_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if user_id not in FULL_OWNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    quizzes = get_quiz_list()
    if not quizzes:
        await update.message.reply_text("📚 لا توجد اختبارات. أرسل ملفاً لإنشاء اختبار.")
        return
    keyboard = []
    for q in quizzes:
        keyboard.append([InlineKeyboardButton(
            f"🗑️ حذف: {q['title'][:25]}", callback_data=f"delete_quiz_{q['id']}"
        )])
    await update.message.reply_text(
        f"📚 *إدارة الاختبارات* ({len(quizzes)} اختبار)\n\nاختر اختباراً للحذف:",
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode='Markdown'
    )

async def handle_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    user_id = query.from_user.id
    data = query.data

    # ── callbacks Google Drive ──
    if data.startswith("drv_"):
        await query.answer()
        await handle_drive_callback(query, user_id, data)
        return

    # ── callbacks منصة التحصيلي ──
    if data.startswith(("pq_", "pqa_", "live_", "asmnt_", "grd_")):
        await query.answer()
        await handle_platform_callback(query, user_id, data)
        return

    await query.answer()

    if data == "show_quizzes":
        quizzes = get_quiz_list()
        if not quizzes:
            await query.message.reply_text("📚 لا توجد اختبارات متاحة.")
            return
        keyboard = [[InlineKeyboardButton(
            f"📝 {q['title']} ({q['questions_count']} سؤال)",
            callback_data=f"start_quiz_{q['id']}"
        )] for q in quizzes]
        await query.message.reply_text(
            f"📚 *الاختبارات المتاحة*\n\n👥 عدد المستخدمين: *{user_count}*",
            reply_markup=InlineKeyboardMarkup(keyboard),
            parse_mode='Markdown'
        )

    elif data.startswith("start_quiz_"):
        quiz_id = data[len("start_quiz_"):]
        quiz_data = load_quiz(quiz_id)
        if not quiz_data:
            await query.message.reply_text("❌ لم يتم العثور على الاختبار.")
            return
        quiz_sessions[user_id] = {'quiz_id': quiz_id, 'current_q': 0, 'score': 0, 'answers': []}
        await query.message.reply_text(
            f"🎯 *{quiz_data['title']}*\n\n"
            f"📊 عدد الأسئلة: {len(quiz_data['questions'])}\n\n"
            "سيتم عرض الأسئلة واحداً تلو الآخر. حظاً موفقاً! 🍀",
            parse_mode='Markdown'
        )
        await send_question(query.message, context, user_id, quiz_data, 0)

    elif data.startswith("answer_"):
        parts = data.split("_")
        q_index = int(parts[1])
        answer_index = int(parts[2])

        session = quiz_sessions.get(user_id)
        if not session:
            await query.message.reply_text("❌ جلسة الاختبار انتهت. استخدم /quiz للبدء.")
            return

        quiz_data = load_quiz(session['quiz_id'])
        if not quiz_data:
            await query.message.reply_text("❌ خطأ في تحميل الاختبار.")
            return

        question = quiz_data['questions'][q_index]
        selected_option = question['options'][answer_index]
        correct_option = question['correct']

        is_correct = selected_option.startswith(correct_option[0]) or selected_option == correct_option
        if is_correct:
            session['score'] += 1
            result_line = "✅ *إجابة صحيحة!*"
        else:
            result_line = f"❌ *إجابة خاطئة!*\n✅ الصحيح: {correct_option}"

        feedback = (
            f"{result_line}\n\n"
            f"💡 *الشرح:* {question.get('explanation', 'لا يوجد شرح')}\n\n"
            f"📊 النتيجة حتى الآن: {session['score']}/{q_index + 1}"
        )
        await query.message.reply_text(feedback, parse_mode='Markdown')

        next_q = q_index + 1
        session['current_q'] = next_q

        if next_q >= len(quiz_data['questions']):
            await show_quiz_result(query.message, context, user_id, quiz_data)
        else:
            await send_question(query.message, context, user_id, quiz_data, next_q)

    elif data.startswith("delete_quiz_") and user_id in FULL_OWNERS:
        quiz_id = data[len("delete_quiz_"):]
        path = os.path.join(QUIZZES_DIR, f"{quiz_id}.json")
        if os.path.exists(path):
            os.remove(path)
            await query.message.reply_text("✅ تم حذف الاختبار.")
        else:
            await query.message.reply_text("❌ لم يتم العثور على الاختبار.")

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if user_id not in FULL_OWNERS:
        await update.message.reply_text("📚 أرسل `/quiz` لحل الاختبارات المتاحة!")
        return

    doc = update.message.document
    file_name = doc.file_name or "file"
    caption = update.message.caption or ""

    await update.message.reply_text("⏳ جاري تنزيل الملف وقراءته...")

    file = await context.bot.get_file(doc.file_id)
    file_path = f"/tmp/{doc.file_id}_{file_name}"
    await file.download_to_drive(file_path)

    ext = file_name.lower().rsplit('.', 1)[-1] if '.' in file_name else ''
    text_content = None
    image_base64 = None

    try:
        if ext == 'pdf':
            text_content = extract_pdf_text(file_path)
            if not text_content:
                await update.message.reply_text("❌ تعذر قراءة PDF.")
                return
        elif ext in ('txt', 'md', 'text'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
        elif ext in ('pptx', 'ppt'):
            text_content = extract_pptx_text(file_path)
            if not text_content:
                await update.message.reply_text("❌ تعذر قراءة PPTX.")
                return
        elif ext in ('jpg', 'jpeg', 'png', 'webp'):
            with open(file_path, 'rb') as f:
                image_base64 = base64.b64encode(f.read()).decode()
        else:
            await update.message.reply_text(
                f"❌ نوع الملف غير مدعوم: .{ext}\n"
                "الأنواع المدعومة: PDF, TXT, PPTX, JPG, PNG"
            )
            return

        if caption:
            await create_quiz_from_content(update, context, text_content, image_base64, caption)
        else:
            pending_quiz[user_id] = {
                'text_content': text_content,
                'image_base64': image_base64
            }
            await update.message.reply_text(
                "✅ *تم استلام الملف!*\n\n"
                "الآن أرسل تعليماتك لإنشاء الاختبار، مثال:\n"
                "• أنشئ 10 أسئلة اختيار من متعدد\n"
                "• أنشئ 5 أسئلة صح وخطأ و5 اختيار متعدد\n"
                "• أنشئ 15 سؤالاً متنوعاً مع التركيز على المفاهيم الأساسية",
                parse_mode='Markdown'
            )
    except Exception as e:
        logging.error(f"خطأ في معالجة الملف: {e}")
        await update.message.reply_text(f"❌ خطأ: {str(e)[:200]}")
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)

async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if user_id not in FULL_OWNERS:
        await update.message.reply_text(
            "📚 لإرسال الصور، تواصل مع المالك.\n"
            "استخدم /quiz لحل الاختبارات المتاحة!"
        )
        return

    caption = update.message.caption or ""
    await update.message.reply_text("⏳ جاري تحليل الصورة...")

    photo = update.message.photo[-1]
    file = await context.bot.get_file(photo.file_id)
    file_path = f"/tmp/{photo.file_id}.jpg"
    await file.download_to_drive(file_path)

    try:
        with open(file_path, 'rb') as f:
            image_base64 = base64.b64encode(f.read()).decode()

        if caption:
            await create_quiz_from_content(update, context, None, image_base64, caption)
        else:
            pending_quiz[user_id] = {'text_content': None, 'image_base64': image_base64}
            await update.message.reply_text(
                "📷 *تم استلام الصورة!*\n\nأرسل تعليماتك لإنشاء الاختبار:",
                parse_mode='Markdown'
            )
    except Exception as e:
        await update.message.reply_text(f"❌ خطأ: {str(e)[:200]}")
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    user_message = update.message.text
    user_name = update.message.from_user.first_name
    username = update.message.from_user.username

    track_user(user_id)
    performance_stats['total_messages_processed'] += 1

    if user_id not in users_db:
        users_db[user_id] = {
            'name': user_name,
            'username': username,
            'first_seen': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }

    if user_id in FULL_OWNERS and user_id in pending_quiz:
        pending = pending_quiz.pop(user_id)
        await create_quiz_from_content(
            update, context,
            pending.get('text_content'),
            pending.get('image_base64'),
            user_message
        )
        return

    if user_id in drive_sessions and drive_sessions[user_id].get('step') == 'ask_count':
        session = drive_sessions.pop(user_id)
        try:
            count = int(re.search(r'\d+', user_message).group())
            count = max(5, min(count, 200))
        except Exception:
            await update.message.reply_text("❌ أرسل رقماً صحيحاً مثل: 20")
            drive_sessions[user_id] = session
            return
        await process_drive_file(update, context, user_id, session['file_id'], session['file_name'], count)
        return

    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action='typing')

    try:
        messages = [{"role": "system", "content": "أنت مساعد ذكي متخصص في الاختبارات الجامعية. رد بالعربية. لديك ذاكرة كاملة للمحادثة."}]
        for old_msg in conversation_history[user_id]:
            messages.append(old_msg)
        messages.append({"role": "user", "content": user_message})
        conversation_history[user_id].append({"role": "user", "content": user_message})

        performance_stats['total_api_calls'] += 1
        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {OPENROUTER_KEY}", "Content-Type": "application/json"},
            json={"model": "meta-llama/llama-3-8b-instruct", "messages": messages, "temperature": 0.3},
            timeout=60
        )
        data = response.json()

        if response.status_code == 200:
            reply = data['choices'][0]['message']['content']
            await update.message.reply_text(reply)
            conversation_history[user_id].append({"role": "assistant", "content": reply})
            if performance_stats['total_messages_processed'] % 20 == 0:
                save_conversations()
            if user_id not in FULL_OWNERS:
                user_info = f"@{username}" if username else f"معرف {user_id}"
                special_name = "فاطمة المطيري" if user_id == FATIMA_ID else ("عبدالخالق" if user_id == ABDULKHALIQ_ID else user_name)
                await send_to_owner(
                    context,
                    f"📩 *محادثة*\n👤 {special_name} | {user_info}\n"
                    f"💬 *السؤال:* {user_message[:150]}\n"
                    f"🤖 *الإجابة:* {reply[:150]}..."
                )
        else:
            error_msg = data.get('error', {}).get('message', 'خطأ غير معروف')
            await update.message.reply_text(f"❌ خطأ: {error_msg}")

    except Exception as e:
        logging.error(f"خطأ: {e}")
        await update.message.reply_text(f"❌ حدث خطأ: {str(e)[:100]}")

# ========== Owner Admin Commands ==========
async def memory_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in FULL_OWNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    memory = get_memory_usage()
    conv_stats = estimate_conversation_size()
    uptime = str(datetime.now() - performance_stats['start_time']).split('.')[0]
    process = psutil.Process(os.getpid())
    system_memory = psutil.virtual_memory()
    text = (
        f"*📊 مراقبة الذاكرة*\n\n"
        f"⏱️ *وقت التشغيل:* {uptime}\n"
        f"🔄 *الرسائل المعالجة:* {performance_stats['total_messages_processed']:,}\n"
        f"🌐 *استدعاءات API:* {performance_stats['total_api_calls']:,}\n\n"
        f"*💾 ذاكرة البوت:*\n"
        f"• المستخدمة: `{memory['rss_human']}`\n"
        f"• نسبة الاستخدام: `{memory['percent']:.2f}%`\n"
        f"• الذروة: `{humanize.naturalsize(performance_stats['peak_memory'])}`\n\n"
        f"*📁 بيانات المحادثات:*\n"
        f"• المستخدمين: `{conv_stats['users_count']}`\n"
        f"• الرسائل: `{conv_stats['total_messages']:,}`\n"
        f"• حجم تقريبي: `{conv_stats['estimated_human']}`\n\n"
        f"*🖥️ النظام:*\n"
        f"• الذاكرة الكلية: `{humanize.naturalsize(system_memory.total)}`\n"
        f"• المتاحة: `{humanize.naturalsize(system_memory.available)}`\n"
        f"• نسبة النظام: `{system_memory.percent}%`"
    )
    await update.message.reply_text(text, parse_mode='Markdown')

async def users_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in FULL_OWNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    if not users_db:
        await update.message.reply_text("📭 لا يوجد مستخدمين حتى الآن.")
        return
    memory = get_memory_usage()
    message = f"*📋 قائمة المستخدمين* (الذاكرة: {memory['rss_human']})\n\n"
    users_list = sorted(users_db.items(), key=lambda x: len(conversation_history.get(x[0], [])), reverse=True)
    for uid, info in users_list:
        conv_length = len(conversation_history.get(uid, []))
        uname = f"@{info['username']}" if info['username'] else "لا يوجد"
        message += f"• *{info['name']}* | 🆔 `{uid}`\n  📱 {uname} | 💬 {conv_length} رسالة | 🕐 {info['first_seen']}\n\n"
        if len(message) > 3500:
            await update.message.reply_text(message, parse_mode='Markdown')
            message = ""
    if message:
        await update.message.reply_text(message, parse_mode='Markdown')

async def stats_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in FULL_OWNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    memory = get_memory_usage()
    conv_stats = estimate_conversation_size()
    uptime = datetime.now() - performance_stats['start_time']
    uptime_str = str(uptime).split('.')[0]
    quizzes = get_quiz_list()
    text = (
        f"*📊 إحصائيات البوت*\n\n"
        f"⏱️ *مدة التشغيل:* {uptime_str}\n"
        f"👥 *المستخدمين:* {user_count}\n"
        f"💬 *الرسائل:* {conv_stats['total_messages']:,}\n"
        f"📚 *الاختبارات:* {len(quizzes)}\n\n"
        f"*💾 الموارد:*\n"
        f"• ذاكرة: {memory['rss_human']} ({memory['percent']:.1f}%)\n"
        f"• API calls: {performance_stats['total_api_calls']:,}\n"
        f"• معدل: {performance_stats['total_messages_processed'] / max(1, uptime.total_seconds() / 3600):.1f} رسالة/ساعة"
    )
    await update.message.reply_text(text, parse_mode='Markdown')

async def send_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in FULL_OWNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    if len(context.args) < 2:
        await update.message.reply_text("❌ الاستخدام: `/send 123456789 الرسالة`", parse_mode='Markdown')
        return
    try:
        target_id = int(context.args[0])
        msg = ' '.join(context.args[1:])
        success = await send_to_user(context, target_id, msg)
        if success:
            await update.message.reply_text(f"✅ تم الإرسال إلى `{target_id}`", parse_mode='Markdown')
            conversation_history[target_id].append({"role": "assistant", "content": msg})
        else:
            await update.message.reply_text(f"❌ فشل الإرسال إلى `{target_id}`", parse_mode='Markdown')
    except ValueError:
        await update.message.reply_text("❌ معرف المستخدم يجب أن يكون رقماً!")
    except Exception as e:
        await update.message.reply_text(f"❌ خطأ: {str(e)}")

async def broadcast_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in FULL_OWNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    if not context.args:
        await update.message.reply_text("❌ الاستخدام: `/broadcast الرسالة`", parse_mode='Markdown')
        return
    msg = ' '.join(context.args)
    if not users_db:
        await update.message.reply_text("📭 لا يوجد مستخدمين.")
        return
    targets = [uid for uid in users_db if uid not in FULL_OWNERS]
    await update.message.reply_text(f"📤 جاري الإرسال إلى {len(targets)} مستخدم...")
    success_count = 0
    fail_count = 0
    for uid in targets:
        if await send_to_user(context, uid, msg):
            success_count += 1
            conversation_history[uid].append({"role": "assistant", "content": msg})
        else:
            fail_count += 1
    await update.message.reply_text(f"✅ تم الإرسال: {success_count}\n❌ فشل: {fail_count}")

async def clear_memory_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if user_id in conversation_history:
        old_size = sum(len(msg.get('content', '')) for msg in conversation_history[user_id]) * 2
        conversation_history[user_id] = []
        await update.message.reply_text(
            f"🧹 *تم مسح ذاكرتك!*\n"
            f"تم تحرير حوالي {humanize.naturalsize(old_size)}",
            parse_mode='Markdown'
        )
    else:
        await update.message.reply_text("لا يوجد محادثات سابقة لمسحها.")

async def clear_all_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in FULL_OWNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    conv_stats = estimate_conversation_size()
    conversation_history.clear()
    performance_stats['total_messages_processed'] = 0
    await update.message.reply_text(
        f"🧹 *تم مسح ذاكرة الجميع!*\n"
        f"تم حذف {conv_stats['total_messages']:,} رسالة",
        parse_mode='Markdown'
    )

async def save_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in FULL_OWNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    if save_conversations():
        file_size = os.path.getsize(HISTORY_FILE) if os.path.exists(HISTORY_FILE) else 0
        await update.message.reply_text(
            f"💾 *تم حفظ المحادثات!*\nحجم الملف: {humanize.naturalsize(file_size)}",
            parse_mode='Markdown'
        )
    else:
        await update.message.reply_text("❌ فشل حفظ المحادثات!")

# ═══════════════════════════════════════════════════════
#        منصة رواد التحصيلي — قراءة الأسئلة والاختبارات
# ═══════════════════════════════════════════════════════

import sqlite3

PLATFORM_DB  = os.path.join(os.path.dirname(__file__), 'coach_tahseeli', 'tahseeli.db')
PLATFORM_URL = os.environ.get('PLATFORM_API_URL', '').rstrip('/')   # مثال: https://rowadaltahseeli1.onrender.com
BOT_API_KEY  = os.environ.get('BOT_API_KEY', '')


def _api_get(endpoint: str, params: dict = None) -> list | dict | None:
    """استدعاء API المنصة — يُعيد None عند الفشل"""
    if not PLATFORM_URL:
        return None
    try:
        headers = {}
        if BOT_API_KEY:
            headers['X-Bot-Key'] = BOT_API_KEY
        url = f"{PLATFORM_URL}/api/bot/{endpoint}"
        resp = requests.get(url, params=params or {}, headers=headers, timeout=15)
        if resp.status_code == 200:
            return resp.json()
        logging.warning(f"[API] {endpoint} → {resp.status_code}")
        return None
    except Exception as e:
        logging.warning(f"[API] فشل الاتصال بالمنصة: {e}")
        return None


SUBJECT_AR = {
    'physics':   'الفيزياء ⚛️',
    'chemistry': 'الكيمياء 🧪',
    'biology':   'الأحياء 🧬',
    'math':      'الرياضيات 📐',
}

EXAM_TYPE_AR = {
    'daily_train': 'تدريب يومي',
    'level_test':  'تقييم المستوى',
    'chapter_1':   'الفصل الأول',
    'chapter_2':   'الفصل الثاني',
    'quick_test':  'اختبار سريع',
    'final_test':  'اختبار شامل',
}

def get_platform_questions(subject: str = None, exam_type: str = None,
                           limit: int = 10, source: str = None,
                           grade: str = None) -> list:
    """قراءة أسئلة المنصة — عبر API أولاً، ثم SQLite محلياً كاحتياطي"""
    # ── API ──────────────────────────────────────────────
    params = {'limit': limit}
    if subject:   params['subject']   = subject
    if exam_type: params['exam_type'] = exam_type
    if source:    params['source']    = source
    if grade:     params['grade']     = grade
    result = _api_get('questions', params)
    if result is not None:
        return result

    # ── SQLite محلي (احتياطي) ────────────────────────────
    if not os.path.exists(PLATFORM_DB):
        return []
    try:
        conn = sqlite3.connect(PLATFORM_DB)
        conn.row_factory = sqlite3.Row
        cur = conn.cursor()

        def _fetch(subj, etype, src=None):
            q = "SELECT * FROM questions WHERE is_active=1"
            p = []
            if subj:
                q += " AND subject=?"; p.append(subj)
            if grade:
                q += " AND grade=?"; p.append(str(grade))
            if src:
                if src == 'prev':
                    q += " AND source LIKE 'prev_%'"
                else:
                    q += " AND source=?"; p.append(src)
            elif etype:
                if etype in ('daily_train', 'level_test', 'quick_test'):
                    pass
                elif etype == 'final_test':
                    q += " AND exam_type IN (?,?,?)"; p.extend(['final_test', 'chapter_1', 'chapter_2'])
                else:
                    q += " AND exam_type=?"; p.append(etype)
            q += " ORDER BY RANDOM() LIMIT ?"; p.append(limit)
            cur.execute(q, p)
            return [dict(r) for r in cur.fetchall()]

        rows = _fetch(subject, exam_type, source)
        if not rows and source:
            rows = _fetch(subject, exam_type)
        conn.close()
        return rows
    except Exception as e:
        logging.error(f"خطأ قاعدة بيانات المنصة: {e}")
        return []


def get_subject_grades(subject: str) -> list:
    """إرجاع قائمة الصفوف المتاحة لمادة معينة"""
    result = _api_get('grades', {'subject': subject})
    if result is not None:
        return result
    if not os.path.exists(PLATFORM_DB):
        return []
    try:
        conn = sqlite3.connect(PLATFORM_DB)
        cur = conn.cursor()
        cur.execute(
            "SELECT DISTINCT grade FROM questions WHERE subject=? AND is_active=1 AND grade IS NOT NULL ORDER BY grade",
            (subject,)
        )
        grades = [r[0] for r in cur.fetchall() if r[0]]
        conn.close()
        return grades
    except:
        return []


GRADE_AR = {'1': 'أول ثانوي', '2': 'ثاني ثانوي', '3': 'ثالث ثانوي'}
GRADE_EMOJI = {'1': '1️⃣', '2': '2️⃣', '3': '3️⃣'}


async def _show_grade_selection(message, subject: str, cmd: str):
    """عرض اختيار الصف الدراسي"""
    grades = get_subject_grades(subject)
    subj_ar = SUBJECT_AR.get(subject, subject)
    if not grades:
        await message.reply_text(f"❌ لا يوجد بنك أسئلة لـ {subj_ar} بعد.")
        return False
    if len(grades) == 1:
        return grades[0]
    keyboard = []
    row = []
    for g in grades:
        label = f"{GRADE_EMOJI.get(g,'🔹')} {GRADE_AR.get(g, f'الصف {g}')}"
        row.append(InlineKeyboardButton(label, callback_data=f"grd_{cmd}_{subject}_{g}"))
        if len(row) == 2:
            keyboard.append(row); row = []
    if row:
        keyboard.append(row)
    await message.reply_text(
        f"📚 *{subj_ar}* — اختر الصف الدراسي:",
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode='Markdown'
    )
    return None

def get_platform_stats() -> dict:
    """إحصائيات منصة رواد التحصيلي"""
    result = _api_get('stats')
    if result is not None:
        return result
    if not os.path.exists(PLATFORM_DB):
        return {}
    try:
        conn = sqlite3.connect(PLATFORM_DB)
        cur = conn.cursor()
        cur.execute("SELECT COUNT(*) FROM questions WHERE is_active=1")
        q_count = cur.fetchone()[0]
        cur.execute("SELECT COUNT(*) FROM users WHERE is_admin=0")
        u_count = cur.fetchone()[0]
        cur.execute("SELECT COUNT(*) FROM evaluations")
        e_count = cur.fetchone()[0]
        conn.close()
        return {'questions': q_count, 'students': u_count, 'exams': e_count}
    except Exception as e:
        logging.error(f"إحصائيات المنصة: {e}")
        return {}

# ── ذكاء اصطناعي مجاني (بدون تكلفة) ──────────────────
FREE_AI_MODELS = [
    'meta-llama/llama-3-8b-instruct:free',
    'google/gemma-2-9b-it:free',
    'mistralai/mistral-7b-instruct:free',
]

def call_free_ai(prompt: str) -> str:
    """استدعاء نموذج ذكاء اصطناعي مجاني"""
    for model in FREE_AI_MODELS:
        try:
            r = requests.post(
                'https://openrouter.ai/api/v1/chat/completions',
                headers={'Authorization': f'Bearer {OPENROUTER_KEY}',
                         'Content-Type': 'application/json'},
                json={'model': model,
                      'messages': [{'role': 'user', 'content': prompt}],
                      'temperature': 0.4, 'max_tokens': 800},
                timeout=30
            )
            data = r.json()
            if r.status_code == 200:
                return data['choices'][0]['message']['content']
        except Exception as e:
            logging.warning(f"Free AI {model}: {e}")
            continue
    return ''

def ai_analyze_result(subject: str, correct: int, total: int) -> str:
    """تحليل نتيجة الاختبار بالذكاء الاصطناعي المجاني"""
    pct = round((correct / total) * 100) if total else 0
    subject_ar = SUBJECT_AR.get(subject, subject)
    prompt = f"""أنت مساعد تعليمي للتحصيل الدراسي السعودي.
الطالب حل اختبار {subject_ar}: أجاب صح على {correct} من {total} ({pct}%).
قدم تحليلاً موجزاً (3 جمل) بالعربية مع توصية واحدة عملية.
أجب مباشرة بدون مقدمة."""
    result = call_free_ai(prompt)
    if not result:
        if pct >= 80:
            return f"أداء ممتاز ✨ حافظ على هذا المستوى وراجع المفاهيم الصعبة."
        elif pct >= 60:
            return f"أداء جيد 👍 راجع الأسئلة التي أخطأت فيها وركّز على فهم المفاهيم الأساسية."
        else:
            return f"تحتاج مزيداً من التدريب 💪 راجع الكتاب المدرسي وحل تمارين إضافية يومياً."
    return result

# ── جلسات الاختبار المباشر في المجموعات ──────────────
# {chat_id: {host_id, subject, questions, current_q, scores:{user_id:{name,score,answered}}, active}}
live_sessions = {}

# ── جلسات اختبار المنصة للمالك في الخاص ──────────────
# {user_id: {subject, exam_type, questions, current_q, score}}
platform_quiz_sessions = {}

# ── أوامر منصة رواد التحصيلي (للمالك) ──────────────────

async def platform_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """عرض قائمة منصة رواد التحصيلي"""
    if update.effective_user.id not in FULL_OWNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    stats = get_platform_stats()
    text = (
        "🏫 *منصة رواد التحصيلي*\n"
        "━━━━━━━━━━━━━━━━━━━━\n"
        f"📚 الأسئلة: *{stats.get('questions', 0)}*\n"
        f"👥 الطلاب: *{stats.get('students', 0)}*\n"
        f"📝 الاختبارات: *{stats.get('exams', 0)}*\n"
        "━━━━━━━━━━━━━━━━━━━━\n\n"
        "*أوامر اختبارات المنصة:*\n"
        "• `/daily physics` — تدريب يومي فيزياء\n"
        "• `/daily chemistry` — تدريب يومي كيمياء\n"
        "• `/daily biology` — تدريب يومي أحياء\n"
        "• `/daily math` — تدريب يومي رياضيات\n\n"
        "• `/assessment physics` — تقييم مستوى\n"
        "• `/ch1 physics` — اختبار الفصل الأول\n"
        "• `/ch2 physics` — اختبار الفصل الثاني\n\n"
        "*الاختبار المباشر في المجموعة:*\n"
        "• `/live physics 10` — ابدأ اختباراً حياً\n"
        "• `/endlive` — أنهِ الاختبار الحالي\n"
        "• `/leaderboard` — نتائج الاختبار"
    )
    keyboard = [[
        InlineKeyboardButton("⚛️ فيزياء", callback_data="pq_daily_physics"),
        InlineKeyboardButton("🧪 كيمياء", callback_data="pq_daily_chemistry"),
    ],[
        InlineKeyboardButton("🧬 أحياء", callback_data="pq_daily_biology"),
        InlineKeyboardButton("📐 رياضيات", callback_data="pq_daily_math"),
    ]]
    await update.message.reply_text(text, parse_mode='Markdown',
                                    reply_markup=InlineKeyboardMarkup(keyboard))

async def _show_assessment_count(message, subject: str, grade: str):
    """عرض اختيار عدد أسئلة التقييم بعد اختيار الصف"""
    import sqlite3 as _sq
    subj_ar = SUBJECT_AR.get(subject, subject)
    grade_ar = GRADE_AR.get(str(grade), f'الصف {grade}')
    try:
        # جرّب API أولاً
        cnt = _api_get('count', {'subject': subject, 'grade': str(grade)})
        if cnt is not None:
            total = cnt.get('count', 0)
        else:
            _conn = _sq.connect(PLATFORM_DB)
            _cur = _conn.cursor()
            _cur.execute(
                "SELECT COUNT(*) FROM questions WHERE subject=? AND grade=? AND is_active=1",
                (subject, str(grade))
            )
            total = _cur.fetchone()[0]
            _conn.close()
    except:
        total = 0
    if total == 0:
        await message.reply_text(f"❌ لا يوجد بنك أسئلة لـ {subj_ar} ({grade_ar}) بعد.")
        return
    keyboard = [
        [
            InlineKeyboardButton("🔟 10 أسئلة",    callback_data=f"asmnt_{subject}_{grade}_10"),
            InlineKeyboardButton("3️⃣0️⃣ 30 سؤال",  callback_data=f"asmnt_{subject}_{grade}_30"),
        ],
        [
            InlineKeyboardButton("5️⃣0️⃣ 50 سؤال",  callback_data=f"asmnt_{subject}_{grade}_50"),
            InlineKeyboardButton("💯 100 سؤال",     callback_data=f"asmnt_{subject}_{grade}_100"),
        ],
        [
            InlineKeyboardButton(f"📚 كل البنك ({total} سؤال)", callback_data=f"asmnt_{subject}_{grade}_{total}"),
        ],
    ]
    await message.reply_text(
        f"📊 *اختبار تقييم — {subj_ar}*\n"
        f"🎓 الصف: *{grade_ar}*\n\n"
        f"🗂 البنك يحتوي على *{total} سؤال*\n"
        f"كل مرة ستحصل على أسئلة مختلفة عشوائياً\n\n"
        f"كم سؤالاً تريد الآن؟",
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode='Markdown'
    )


async def _start_platform_quiz(update, context, subject, exam_type, count=10, grade=None):
    """بدء اختبار من المنصة في خاص المالك"""
    user_id = update.effective_user.id
    questions = get_platform_questions(subject, exam_type, count, grade=grade)
    if not questions:
        await update.message.reply_text(
            f"❌ لا توجد أسئلة في قاعدة البيانات للمادة المطلوبة.\n"
            f"أضف أسئلة أولاً من لوحة المشرف على المنصة."
        )
        return
    import random; random.shuffle(questions)
    platform_quiz_sessions[user_id] = {
        'subject': subject, 'exam_type': exam_type,
        'questions': questions, 'current_q': 0, 'score': 0
    }
    subj_ar  = SUBJECT_AR.get(subject, subject)
    etype_ar = EXAM_TYPE_AR.get(exam_type, exam_type)
    grade_ar = f" — {GRADE_AR.get(str(grade), '')}" if grade else ""
    await update.message.reply_text(
        f"🎯 *{subj_ar}{grade_ar} — {etype_ar}*\n"
        f"📊 عدد الأسئلة: {len(questions)}\n\n"
        "سيتم عرض الأسئلة واحداً تلو الآخر. حظاً موفقاً! 🍀",
        parse_mode='Markdown'
    )
    await _send_platform_question(update.message, context, user_id)

async def _send_platform_question(send_target, context, user_id):
    """إرسال سؤال منصة للمالك"""
    session = platform_quiz_sessions.get(user_id)
    if not session:
        return
    qs = session['questions']
    idx = session['current_q']
    if idx >= len(qs):
        await _show_platform_result(send_target, context, user_id)
        return
    q = qs[idx]
    total = len(qs)
    opts = [
        ('A', q['option_a']), ('B', q['option_b']),
        ('C', q['option_c']), ('D', q['option_d'])
    ]
    keyboard = [[InlineKeyboardButton(f"{l}. {v[:40]}", callback_data=f"pqa_{user_id}_{idx}_{l}")]
                for l, v in opts]
    await send_target.reply_text(
        f"❓ *سؤال {idx+1} من {total}*\n\n{q['text']}",
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode='Markdown'
    )

async def _show_platform_result(send_target, context, user_id):
    """عرض نتيجة اختبار المنصة"""
    session = platform_quiz_sessions.pop(user_id, {})
    correct = session.get('score', 0)
    total   = len(session.get('questions', [1]))
    pct     = round((correct / total) * 100) if total else 0
    emoji   = '🏆' if pct >= 90 else '⭐' if pct >= 75 else '👍' if pct >= 60 else '💪'
    subject  = session.get('subject', '')
    analysis = ai_analyze_result(subject, correct, total)
    await send_target.reply_text(
        f"🎯 *انتهى الاختبار!*\n\n"
        f"📊 النتيجة: *{correct}/{total}* ({pct}%) {emoji}\n\n"
        f"🤖 *تحليل الذكاء الاصطناعي:*\n{analysis}",
        parse_mode='Markdown'
    )

async def daily_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """أمر التدريب اليومي"""
    if update.effective_user.id not in EXAM_PARTNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    subject = context.args[0].lower() if context.args else ''
    if subject not in SUBJECT_AR:
        await update.message.reply_text(
            "❌ اختر مادة:\n`/daily physics` | `/daily chemistry` | `/daily biology` | `/daily math`",
            parse_mode='Markdown'
        )
        return
    result = await _show_grade_selection(update.message, subject, 'daily')
    if result is False:
        return
    if result is None:
        return
    await _start_platform_quiz(update, context, subject, 'daily_train', 10, grade=result)

async def assessment_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """أمر اختبار التقييم — يعرض اختيار الصف ثم عدد الأسئلة"""
    if update.effective_user.id not in EXAM_PARTNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    subject = context.args[0].lower() if context.args else ''
    if subject not in SUBJECT_AR:
        await update.message.reply_text(
            "❌ اختر مادة:\n`/assessment physics` | `/assessment chemistry` | `/assessment biology` | `/assessment math`",
            parse_mode='Markdown'
        )
        return
    result = await _show_grade_selection(update.message, subject, 'asmnt')
    if result is False:
        return
    if result is None:
        return
    await _show_assessment_count(update.message, subject, result)

async def ch1_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in EXAM_PARTNERS:
        await update.message.reply_text("❌"); return
    subject = context.args[0].lower() if context.args else ''
    if subject not in SUBJECT_AR:
        await update.message.reply_text("❌ `/ch1 physics`", parse_mode='Markdown'); return
    await _start_platform_quiz(update, context, subject, 'chapter_1', 10)

async def ch2_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in EXAM_PARTNERS:
        await update.message.reply_text("❌"); return
    subject = context.args[0].lower() if context.args else ''
    if subject not in SUBJECT_AR:
        await update.message.reply_text("❌ `/ch2 physics`", parse_mode='Markdown'); return
    await _start_platform_quiz(update, context, subject, 'chapter_2', 10)

# ── الاختبار المباشر في المجموعة ──────────────────────

async def live_quiz_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """بدء اختبار مباشر في المجموعة — للمالك فقط"""
    chat_id = update.effective_chat.id
    user_id = update.effective_user.id
    if user_id not in EXAM_PARTNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    if chat_id == user_id:
        await update.message.reply_text("❌ هذا الأمر للمجموعات فقط. أضف البوت لمجموعتك أولاً.")
        return
    if chat_id in live_sessions and live_sessions[chat_id].get('active'):
        await update.message.reply_text("⚠️ يوجد اختبار مباشر نشط. أنهِه أولاً بـ `/endlive`",
                                        parse_mode='Markdown')
        return
    subject = context.args[0].lower() if context.args else 'physics'
    if subject not in SUBJECT_AR:
        subject = 'physics'
    count = 10
    if len(context.args) >= 2:
        try: count = min(int(context.args[1]), 20)
        except: pass

    questions = get_platform_questions(subject, None, count)
    if not questions:
        await update.message.reply_text("❌ لا توجد أسئلة في قاعدة البيانات بعد!")
        return
    import random; random.shuffle(questions)

    live_sessions[chat_id] = {
        'host_id': user_id, 'subject': subject,
        'questions': questions, 'current_q': 0,
        'scores': {}, 'active': True,
        'answered_current': set()
    }
    subj_ar = SUBJECT_AR.get(subject, subject)
    await update.message.reply_text(
        f"🏆 *اختبار مباشر — {subj_ar}*\n"
        f"━━━━━━━━━━━━━━━━━━━━\n"
        f"📊 عدد الأسئلة: *{len(questions)}*\n"
        f"👥 شارك مع أصدقائك!\n\n"
        f"🔔 سيبدأ الاختبار الآن...",
        parse_mode='Markdown'
    )
    await _send_live_question(context, chat_id)

async def _send_live_question(context, chat_id):
    """إرسال سؤال الاختبار المباشر للمجموعة"""
    session = live_sessions.get(chat_id)
    if not session:
        return
    qs  = session['questions']
    idx = session['current_q']
    if idx >= len(qs):
        await _end_live_quiz(context, chat_id)
        return
    q = qs[idx]
    total = len(qs)
    session['answered_current'] = set()
    opts = [('A', q['option_a']), ('B', q['option_b']),
            ('C', q['option_c']), ('D', q['option_d'])]
    keyboard = [[InlineKeyboardButton(f"{l}. {v[:35]}", callback_data=f"live_{chat_id}_{idx}_{l}")]
                for l, v in opts]
    await context.bot.send_message(
        chat_id=chat_id,
        text=f"❓ *سؤال {idx+1} من {total}*\n\n{q['text']}",
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode='Markdown'
    )

async def end_live_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """إنهاء الاختبار المباشر"""
    chat_id = update.effective_chat.id
    user_id = update.effective_user.id
    if user_id not in EXAM_PARTNERS:
        await update.message.reply_text("❌ هذا الأمر للمالك فقط!")
        return
    if chat_id not in live_sessions:
        await update.message.reply_text("❌ لا يوجد اختبار مباشر نشط.")
        return
    await _end_live_quiz(context, chat_id, force=True)

async def _end_live_quiz(context, chat_id, force=False):
    """إنهاء الاختبار وعرض النتائج"""
    session = live_sessions.pop(chat_id, None)
    if not session:
        return
    scores = session.get('scores', {})
    total_q = session['current_q'] if not force else session['current_q']
    if not scores:
        await context.bot.send_message(
            chat_id=chat_id,
            text="🏁 *انتهى الاختبار المباشر!*\n\nلم يشارك أحد 😔",
            parse_mode='Markdown'
        )
        return
    sorted_scores = sorted(scores.items(), key=lambda x: x[1]['score'], reverse=True)
    medal = ['🥇', '🥈', '🥉']
    board = ""
    for i, (uid, info) in enumerate(sorted_scores):
        m = medal[i] if i < 3 else f"{i+1}."
        pct = round((info['score'] / max(info['answered'], 1)) * 100) if info.get('answered') else 0
        board += f"{m} *{info['name']}* — {info['score']} نقطة ({pct}%)\n"
    await context.bot.send_message(
        chat_id=chat_id,
        text=f"🏆 *نتائج الاختبار المباشر!*\n"
             f"━━━━━━━━━━━━━━━━━━━━\n{board}\n"
             f"شكراً للمشاركة! 🙏",
        parse_mode='Markdown'
    )

async def leaderboard_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """عرض نتائج الاختبار الحالي"""
    chat_id = update.effective_chat.id
    session = live_sessions.get(chat_id)
    if not session:
        await update.message.reply_text("❌ لا يوجد اختبار مباشر نشط حالياً.")
        return
    scores = session.get('scores', {})
    if not scores:
        await update.message.reply_text("📊 لم يجب أحد بعد على الأسئلة.")
        return
    sorted_scores = sorted(scores.items(), key=lambda x: x[1]['score'], reverse=True)
    medal = ['🥇', '🥈', '🥉']
    board = f"📊 *النتائج الحالية — سؤال {session['current_q']}/{len(session['questions'])}*\n━━━━━━━\n"
    for i, (uid, info) in enumerate(sorted_scores):
        m = medal[i] if i < 3 else f"{i+1}."
        board += f"{m} {info['name']} — {info['score']} نقطة\n"
    await update.message.reply_text(board, parse_mode='Markdown')

# ═══════════════════════════════════════════════════════
#        معالج callbacks محدّث يشمل أسئلة المنصة
# ═══════════════════════════════════════════════════════

async def handle_platform_callback(query, user_id, data):
    """معالجة callbacks منصة التحصيلي"""

    # اختيار الصف ثم توجيه للأمر المناسب
    if data.startswith("grd_"):
        # grd_{cmd}_{subject}_{grade}
        parts = data.split("_")
        cmd     = parts[1]
        subject = parts[2]
        grade   = parts[3]
        if cmd == 'asmnt':
            await _show_assessment_count(query.message, subject, grade)
        elif cmd == 'daily':
            import random as _rnd
            questions = get_platform_questions(subject, 'daily_train', 10, grade=grade)
            if not questions:
                questions = get_platform_questions(subject, None, 10, grade=grade)
            if not questions:
                await query.message.reply_text("❌ لا توجد أسئلة بعد.")
                return
            _rnd.shuffle(questions)
            subj_ar  = SUBJECT_AR.get(subject, subject)
            grade_ar = GRADE_AR.get(str(grade), '')
            await query.message.reply_text(
                f"🎯 *تدريب يومي — {subj_ar}*\n"
                f"🎓 الصف: *{grade_ar}*\n"
                f"📊 عدد الأسئلة: {len(questions)}\n\n"
                "حظاً موفقاً! 🍀",
                parse_mode='Markdown'
            )
            platform_quiz_sessions[user_id] = {
                'subject': subject, 'exam_type': 'daily_train',
                'questions': questions, 'current_q': 0, 'score': 0
            }
            await _send_platform_question(query.message, None, user_id)
        elif cmd == 'quick':
            import random as _rnd
            questions = get_platform_questions(subject, 'quick_test', 10, grade=grade)
            if not questions:
                questions = get_platform_questions(subject, None, 10, grade=grade)
            if not questions:
                await query.message.reply_text("❌ لا توجد أسئلة بعد.")
                return
            _rnd.shuffle(questions)
            subj_ar  = SUBJECT_AR.get(subject, subject)
            grade_ar = GRADE_AR.get(str(grade), '')
            await query.message.reply_text(
                f"⚡ *اختبار سريع — {subj_ar}*\n"
                f"🎓 الصف: *{grade_ar}*\n"
                f"📊 عدد الأسئلة: {len(questions)}\n\n"
                "حظاً موفقاً! 🍀",
                parse_mode='Markdown'
            )
            platform_quiz_sessions[user_id] = {
                'subject': subject, 'exam_type': 'quick_test',
                'questions': questions, 'current_q': 0, 'score': 0
            }
            await _send_platform_question(query.message, None, user_id)
        return

    # اختيار عدد أسئلة التقييم من البنك
    if data.startswith("asmnt_"):
        # asmnt_{subject}_{grade}_{count}  أو  asmnt_{subject}_{count} (قديم)
        parts = data.split("_")
        subject = parts[1]
        if len(parts) == 4:
            grade = parts[2]
            count = int(parts[3])
        else:
            grade = None
            count = int(parts[2])
        import random as _rnd
        questions = get_platform_questions(subject, 'level_test', count, grade=grade)
        if not questions:
            await query.message.reply_text("❌ لا توجد أسئلة في البنك حالياً.")
            return
        _rnd.shuffle(questions)
        platform_quiz_sessions[user_id] = {
            'subject': subject, 'exam_type': 'level_test',
            'questions': questions, 'current_q': 0, 'score': 0
        }
        subj_ar  = SUBJECT_AR.get(subject, subject)
        grade_ar = f" — {GRADE_AR.get(str(grade), '')}" if grade else ""
        await query.message.reply_text(
            f"🎯 *اختبار تقييم — {subj_ar}{grade_ar}*\n"
            f"📊 عدد الأسئلة: {len(questions)} (عشوائية من البنك)\n\n"
            "حظاً موفقاً! 🍀",
            parse_mode='Markdown'
        )
        await _send_platform_question(query.message, None, user_id)
        return

    # قائمة المنصة السريعة
    if data.startswith("pq_"):
        parts = data.split("_")
        # pq_daily_physics
        exam_type_raw = parts[1]
        subject       = parts[2]
        # تحويل الاختصارات إلى أنواع قاعدة البيانات الفعلية
        EXAM_TYPE_MAP = {
            'daily':  'daily_train',
            'quick':  'quick_test',
            'level':  'level_test',
            'final':  'final_test',
            'ch1':    'chapter_1',
            'ch2':    'chapter_2',
        }
        exam_type = EXAM_TYPE_MAP.get(exam_type_raw, exam_type_raw)
        await query.message.reply_text(f"⏳ جاري تحميل أسئلة {SUBJECT_AR.get(subject,subject)}...")
        questions = get_platform_questions(subject, exam_type, 10)
        if not questions:
            await query.message.reply_text("❌ لا توجد أسئلة متاحة. أضفها من لوحة المشرف.")
            return
        import random; random.shuffle(questions)
        platform_quiz_sessions[user_id] = {
            'subject': subject, 'exam_type': exam_type,
            'questions': questions, 'current_q': 0, 'score': 0
        }
        await _send_platform_question(query.message, None, user_id)
        return

    # إجابة اختبار المنصة في الخاص
    if data.startswith("pqa_"):
        # pqa_{user_id}_{q_idx}_{letter}
        parts = data.split("_")
        target_uid = int(parts[1])
        q_idx      = int(parts[2])
        letter     = parts[3]
        if user_id != target_uid:
            await query.answer("❌ هذا الاختبار ليس لك!", show_alert=True)
            return
        session = platform_quiz_sessions.get(user_id)
        if not session or session['current_q'] != q_idx:
            await query.answer("⚠️ السؤال انتهى أو لا توجد جلسة.")
            return
        q = session['questions'][q_idx]
        correct = q['answer'].upper()
        is_ok   = letter.upper() == correct
        if is_ok:
            session['score'] += 1
            fb = f"✅ *إجابة صحيحة!*"
        else:
            opt_map = {'A': q['option_a'], 'B': q['option_b'],
                       'C': q['option_c'], 'D': q['option_d']}
            fb = f"❌ *خطأ!*\n✅ الصواب: {correct}. {opt_map.get(correct,'')}"
        explanation = q.get('explanation', '')
        if explanation:
            fb += f"\n\n💡 {explanation}"
        fb += f"\n\n📊 {session['score']}/{q_idx+1}"
        session['current_q'] += 1
        await query.message.reply_text(fb, parse_mode='Markdown')
        await query.answer()
        # السؤال التالي
        if session['current_q'] >= len(session['questions']):
            await _show_platform_result(query.message, None, user_id)
        else:
            await _send_platform_question(query.message, None, user_id)
        return

    # إجابة الاختبار المباشر في المجموعة
    if data.startswith("live_"):
        # live_{chat_id}_{q_idx}_{letter}
        parts   = data.split("_")
        chat_id = int(parts[1])
        q_idx   = int(parts[2])
        letter  = parts[3]
        session = live_sessions.get(chat_id)
        if not session or not session.get('active'):
            await query.answer("❌ الاختبار انتهى.", show_alert=True)
            return
        if session['current_q'] != q_idx:
            await query.answer("⏭️ تأخرت — السؤال انتهى.")
            return
        # منع الإجابة المتكررة
        if user_id in session['answered_current']:
            await query.answer("✋ أجبت بالفعل على هذا السؤال!", show_alert=True)
            return
        session['answered_current'].add(user_id)
        q       = session['questions'][q_idx]
        correct = q['answer'].upper()
        is_ok   = letter.upper() == correct
        # تسجيل النقطة
        name    = query.from_user.first_name or f"مستخدم{user_id}"
        if user_id not in session['scores']:
            session['scores'][user_id] = {'name': name, 'score': 0, 'answered': 0}
        session['scores'][user_id]['answered'] += 1
        if is_ok:
            session['scores'][user_id]['score'] += 1
            await query.answer("✅ إجابة صحيحة! +1", show_alert=False)
        else:
            await query.answer(f"❌ خطأ! الصواب: {correct}", show_alert=False)
        # بعد فترة انتقال إلى السؤال التالي إذا أجاب الكل (اختياري)
        # ننتقل للسؤال التالي فقط عندما يضغط عليه المالك عبر /next (أو تلقائياً بعد X ثانية)
        # لتبسيط التجربة: ننتقل تلقائياً بعد أن يصل عدد الإجابات لعدد معقول
        return

async def next_question_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """ينتقل المالك للسؤال التالي في الاختبار المباشر"""
    chat_id = update.effective_chat.id
    user_id = update.effective_user.id
    if user_id not in EXAM_PARTNERS:
        return
    session = live_sessions.get(chat_id)
    if not session or not session.get('active'):
        await update.message.reply_text("❌ لا يوجد اختبار مباشر.")
        return
    # إظهار نتيجة السؤال الحالي
    q_idx = session['current_q']
    if q_idx < len(session['questions']):
        q       = session['questions'][q_idx]
        correct = q['answer'].upper()
        opt_map = {'A': q['option_a'], 'B': q['option_b'],
                   'C': q['option_c'], 'D': q['option_d']}
        answered_count = len(session['answered_current'])
        await update.message.reply_text(
            f"✅ الإجابة الصحيحة: *{correct}. {opt_map.get(correct, '')}*\n"
            f"👥 عدد المجيبين: {answered_count}\n\n"
            f"💡 {q.get('explanation', '')}",
            parse_mode='Markdown'
        )
    session['current_q'] += 1
    # تحقق من انتهاء الأسئلة
    if session['current_q'] >= len(session['questions']):
        await _end_live_quiz(context, chat_id)
    else:
        await _send_live_question(context, chat_id)

# ── مجموعة المصرّح لهم (المالك + فاطمة) ──
AUTHORIZED_IDS = EXAM_PARTNERS

def _is_authorized(user_id: int) -> bool:
    return user_id in AUTHORIZED_IDS

# ══════════════════════════════════════════════════════════════
#  أوامر الأنواع المختلفة — متاحة للمالك وفاطمة فقط
# ══════════════════════════════════════════════════════════════

async def prev_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """/prev [مادة] — نماذج اختبارات سابقة"""
    user_id = update.effective_user.id
    if not _is_authorized(user_id):
        await update.message.reply_text("❌ هذا الأمر للمالك وفاطمة فقط!")
        return
    args = context.args
    subject = args[0].lower() if args else None
    if subject and subject not in SUBJECT_AR:
        await update.message.reply_text(
            "❌ المادة غير صحيحة.\nالمواد المتاحة: physics | chemistry | biology | math"
        )
        return
    subj_ar = SUBJECT_AR.get(subject, 'جميع المواد')
    intro = (
        f"📋 *نماذج اختبارات سابقة — {subj_ar}*\n"
        "━━━━━━━━━━━━━━━━━━━━\n"
        "📌 *ما هي النماذج السابقة؟*\n"
        "هي اختبارات أُجريت في السنوات الماضية (1443–1446هـ) "
        "وتعكس النمط الحقيقي لأسئلة التحصيلي.\n\n"
        "💡 *لماذا التدرب عليها؟*\n"
        "• تعرّف على أسلوب الصياغة الرسمي\n"
        "• تدرّب على توزيع الوقت\n"
        "• اكتشف المواضيع الأكثر تكراراً\n"
        "━━━━━━━━━━━━━━━━━━━━\n"
        "🚀 يبدأ الاختبار الآن..."
    )
    await update.message.reply_text(intro, parse_mode='Markdown')
    await _start_platform_quiz(update, context, subject, 'final_test', 10)

async def expected_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """/expected [مادة] — الاختبارات السريعة"""
    user_id = update.effective_user.id
    if not _is_authorized(user_id):
        await update.message.reply_text("❌ هذا الأمر للمالك وفاطمة فقط!")
        return
    args = context.args
    subject = args[0].lower() if args else None
    if not subject or subject not in SUBJECT_AR:
        await update.message.reply_text(
            "❌ المادة غير صحيحة.\nالمواد المتاحة: physics | chemistry | biology | math"
        )
        return
    result = await _show_grade_selection(update.message, subject, 'quick')
    if result is False:
        return
    if result is None:
        return
    questions = get_platform_questions(subject, 'quick_test', 10, grade=result)
    if not questions:
        questions = get_platform_questions(subject, None, 10, grade=result)
    if not questions:
        await update.message.reply_text("❌ لا توجد أسئلة بعد.")
        return
    import random; random.shuffle(questions)
    subj_ar = SUBJECT_AR.get(subject, subject)
    grade_ar = GRADE_AR.get(str(result), '')
    await update.message.reply_text(
        f"⚡ *اختبار سريع — {subj_ar}*\n"
        f"🎓 الصف: *{grade_ar}*\n"
        f"📊 عدد الأسئلة: {len(questions)}\n\n"
        "حظاً موفقاً! 🍀",
        parse_mode='Markdown'
    )
    platform_quiz_sessions[user_id] = {
        'subject': subject, 'exam_type': 'quick_test',
        'questions': questions, 'current_q': 0, 'score': 0
    }
    await _send_platform_question(update.message, context, user_id)

async def final_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """/final [مادة] — اختبار شامل لكل الصفوف"""
    user_id = update.effective_user.id
    if not _is_authorized(user_id):
        await update.message.reply_text("❌ هذا الأمر للمالك وفاطمة فقط!")
        return
    args = context.args
    subject = args[0].lower() if args else None
    if subject and subject not in SUBJECT_AR:
        await update.message.reply_text(
            "❌ المادة غير صحيحة.\nالمواد المتاحة: physics | chemistry | biology | math"
        )
        return
    subj_ar = SUBJECT_AR.get(subject, 'جميع المواد')
    intro = (
        f"🌐 *الاختبار الشامل — {subj_ar}*\n"
        "━━━━━━━━━━━━━━━━━━━━\n"
        "📌 *ما هو الاختبار الشامل؟*\n"
        "يغطي منهج الصفوف الثلاثة (الأول والثاني والثالث ثانوي) "
        "في مادة واحدة — تماماً كما في التحصيلي الحقيقي.\n\n"
        "💡 *مناسب لـ:*\n"
        "• من أنهى جميع الفصول ويريد مراجعة شاملة\n"
        "• من يستعد للتحصيلي خلال أسابيع قليلة\n"
        "━━━━━━━━━━━━━━━━━━━━\n"
        "📊 يبدأ الاختبار الشامل الآن..."
    )
    await update.message.reply_text(intro, parse_mode='Markdown')
    await _start_platform_quiz(update, context, subject, 'final_test', 15)

async def arena_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """/arena [مادة] — اختبار ساحة التنافس بين الرواد"""
    user_id = update.effective_user.id
    if not _is_authorized(user_id):
        await update.message.reply_text("❌ هذا الأمر للمالك وفاطمة فقط!")
        return
    args = context.args
    subject = args[0].lower() if args else None
    if subject and subject not in SUBJECT_AR:
        await update.message.reply_text(
            "❌ المادة غير صحيحة.\nالمواد المتاحة: physics | chemistry | biology | math"
        )
        return
    subj_ar = SUBJECT_AR.get(subject, 'جميع المواد')
    intro = (
        f"⚔️ *ساحة التنافس — {subj_ar}*\n"
        "━━━━━━━━━━━━━━━━━━━━\n"
        "🏆 *اختبر نفسك بين الرواد!*\n"
        "أسئلة ساحة التنافس مصممة خصيصاً لقياس مستواك "
        "التنافسي مقارنةً بأقرانك.\n\n"
        "⚡ *مميزات هذا الاختبار:*\n"
        "• أسئلة مختارة بعناية بمستويات متدرجة\n"
        "• تنافس حقيقي — هل أنت من سيحتل القمة؟\n"
        "• كُن رائداً لا تابعاً! 🚀\n"
        "━━━━━━━━━━━━━━━━━━━━\n"
        "⚔️ تبدأ المعركة الآن..."
    )
    await update.message.reply_text(intro, parse_mode='Markdown')
    questions = get_platform_questions(subject, 'final_test', 10, source='arena')
    if not questions:
        questions = get_platform_questions(subject, 'final_test', 10)
    if not questions:
        await update.message.reply_text("❌ لا توجد أسئلة للساحة بعد. أضفها من لوحة المشرف.")
        return
    import random; random.shuffle(questions)
    platform_quiz_sessions[user_id] = {
        'subject': subject, 'exam_type': 'final_test',
        'questions': questions, 'current_q': 0, 'score': 0
    }
    await _send_platform_question(update.message, context, user_id)

async def cmds_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """/cmds — قائمة الأوامر المتاحة للمالك وفاطمة"""
    user_id = update.effective_user.id
    if not _is_authorized(user_id):
        return
    text = (
        "📋 *أوامر منصة رواد التحصيلي*\n"
        "━━━━━━━━━━━━━━━━━━━━\n\n"
        "*📚 اختبارات الصفوف والفصول:*\n"
        "`/ch1 [مادة]` — اختبار الفصل الأول\n"
        "`/ch2 [مادة]` — اختبار الفصل الثاني\n"
        "`/daily [مادة]` — تدريب يومي\n\n"
        "*🌐 الاختبارات الشاملة:*\n"
        "`/final [مادة]` — اختبار شامل كل الصفوف\n\n"
        "*📋 النماذج السابقة:*\n"
        "`/prev [مادة]` — نماذج من السنوات الماضية\n\n"
        "*🔮 الاختبارات المتوقعة:*\n"
        "`/expected [مادة]` — المواضيع المتوقعة\n\n"
        "*⚔️ ساحة التنافس:*\n"
        "`/arena [مادة]` — اختبار ساحة الرواد\n\n"
        "*🎯 التقييم:*\n"
        "`/assessment [مادة]` — تقييم مستوى\n\n"
        "*📡 الاختبار المباشر في المجموعة:*\n"
        "`/live [مادة] [عدد]` — ابدأ اختباراً حياً\n"
        "`/next` — السؤال التالي\n"
        "`/leaderboard` — ترتيب المتنافسين\n"
        "`/endlive` — إنهاء الاختبار\n\n"
        "المواد: `physics` | `chemistry` | `biology` | `math`"
    )
    await update.message.reply_text(text, parse_mode='Markdown')


# ========== Error Handler ==========
async def error_handler(update, context: ContextTypes.DEFAULT_TYPE):
    error = context.error
    if "Conflict" in str(error):
        logging.critical("❌ تعارض: نسخة أخرى من البوت تعمل! أوقف النشر القديم أولاً.")
    else:
        logging.error(f"خطأ في البوت: {error}")

# ========== Main ==========
def main():
    app = Application.builder().token(BOT_TOKEN).build()
    app.add_error_handler(error_handler)

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("quiz", quiz_command))
    app.add_handler(CommandHandler("quizzes", quizzes_manage_command))
    app.add_handler(CommandHandler("users", users_command))
    app.add_handler(CommandHandler("stats", stats_command))
    app.add_handler(CommandHandler("memory", memory_command))
    app.add_handler(CommandHandler("send", send_command))
    app.add_handler(CommandHandler("broadcast", broadcast_command))
    app.add_handler(CommandHandler("clear", clear_memory_command))
    app.add_handler(CommandHandler("clear_all", clear_all_command))
    app.add_handler(CommandHandler("save", save_command))
    # ── منصة رواد التحصيلي ──
    app.add_handler(CommandHandler("platform",   platform_command))
    app.add_handler(CommandHandler("daily",      daily_command))
    app.add_handler(CommandHandler("assessment", assessment_command))
    app.add_handler(CommandHandler("ch1",        ch1_command))
    app.add_handler(CommandHandler("ch2",        ch2_command))
    # ── أوامر الأنواع (مالك + فاطمة) ──
    app.add_handler(CommandHandler("prev",       prev_command))
    app.add_handler(CommandHandler("expected",   expected_command))
    app.add_handler(CommandHandler("final",      final_command))
    app.add_handler(CommandHandler("arena",      arena_command))
    app.add_handler(CommandHandler("cmds",       cmds_command))
    # ── كتب Google Drive ──
    app.add_handler(CommandHandler("book",        drive_book_command))
    # ── الاختبار المباشر ──
    app.add_handler(CommandHandler("live",        live_quiz_command))
    app.add_handler(CommandHandler("endlive",     end_live_command))
    app.add_handler(CommandHandler("next",        next_question_command))
    app.add_handler(CommandHandler("leaderboard", leaderboard_command))

    app.add_handler(CallbackQueryHandler(handle_callback))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))

    app.job_queue.run_repeating(update_bot_description, interval=300, first=5)

    memory = get_memory_usage()
    conv_stats = estimate_conversation_size()
    quizzes = get_quiz_list()
    print("=" * 60)
    print("🤖 بوت الاختبارات الجامعية بالذكاء الاصطناعي")
    print("=" * 60)
    print(f"💾 الذاكرة: {memory['rss_human']} ({memory['percent']:.1f}%)")
    print(f"💬 الرسائل المحفوظة: {conv_stats['total_messages']:,}")
    print(f"👥 المستخدمين: {user_count}")
    print(f"📚 الاختبارات: {len(quizzes)}")
    print("=" * 60)

    try:
        app.run_polling(drop_pending_updates=True, allowed_updates=Update.ALL_TYPES)
    finally:
        save_conversations()
        print("💾 تم حفظ المحادثات.")

if __name__ == '__main__':
    main()
